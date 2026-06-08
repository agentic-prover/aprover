#!/usr/bin/env python3
"""Convert the animated HTML deck (bmc-agent-onepager.html) into a multi-slide
PPTX — one full-bleed 16:9 slide per scene — with Fade animations between slides,
for inserting into an existing deck.

  1) render scenes:  node _shoot_scenes.mjs   (writes _slides/scene_*.png)
  2) python3 build_deck_pptx.py                -> aprover-deck.pptx
"""
import os, glob
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
SLIDES = sorted(glob.glob(os.path.join(HERE, "_slides", "scene_*.png")))
OUT = os.path.join(HERE, "aprover-deck.pptx")

def add_fade_transition(slide, speed="med"):
    """Inject a legacy <p:transition><p:fade/></p:transition> so each slide
    animates (fades) on advance. Placed after cSld/clrMapOvr per CT_Slide order."""
    sld = slide._element
    trans = etree.SubElement(sld, qn("p:transition"))
    trans.set("spd", speed)
    etree.SubElement(trans, qn("p:fade"))
    anchor = sld.find(qn("p:clrMapOvr"))
    if anchor is None:
        anchor = sld.find(qn("p:cSld"))
    anchor.addnext(trans)

prs = Presentation()
prs.slide_width = Inches(13.333)      # 16:9 widescreen (PowerPoint default)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for png in SLIDES:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_fade_transition(s)
prs.save(OUT)
print(f"wrote {OUT} — {len(SLIDES)} slides (16:9, full-bleed, Fade transitions)")
