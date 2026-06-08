#!/usr/bin/env python3
"""Build a single-slide PPTX presenting agentic-prover, with the bmc-agent demo.

  python3 presentation/build_ppt.py

Output: presentation/agentic-prover.pptx
- If presentation/bmc-agent-demo.mp4 exists, it is EMBEDDED as a playable movie.
- Otherwise a generated poster (presentation/_demo_poster.png) is placed and
  click-linked to the live HTML demo (bmc-agent.html).
"""
import os
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
MP4 = os.path.join(HERE, "bmc-agent-demo.mp4")
POSTER = os.path.join(HERE, "_demo_poster.png")
OUT = os.path.join(HERE, "agentic-prover.pptx")

# ---- palette ----
def C(h): return RGBColor.from_string(h)
BG, BG2 = "0d1117", "05070b"
PANEL, PANEL2, LINE = "161b22", "1c2230", "30363d"
TEXT, MUTED = "e6edf3", "8b949e"
GREEN, BLUE, VIOLET, RED, AMBER, GREY = "3fb950", "58a6ff", "a371f7", "f85149", "d29922", "888888"

# ===================================================================
# 1) generate the demo poster (looks like a frame of the animation)
# ===================================================================
def make_poster():
    W, H = 1600, 900
    img = Image.new("RGB", (W, H), (5, 7, 11))
    d = ImageDraw.Draw(img, "RGBA")
    def rgb(h): return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    # module wall
    COLS, ROWS = 8, 5
    aw, ah = 1180, 520
    ax, ay = (W - aw) // 2, (H - ah) // 2 + 20
    gap = 14
    cw = (aw - (COLS - 1) * gap) / COLS
    ch = (ah - (ROWS - 1) * gap) / ROWS
    def cell(i):
        c, r = i % COLS, i // COLS
        x = ax + c * (cw + gap); y = ay + r * (ch + gap)
        return x, y, x + cw, y + ch
    def center(i):
        x0, y0, x1, y1 = cell(i); return ((x0 + x1) / 2, (y0 + y1) / 2)

    bug = 19  # col3,row2
    for i in range(COLS * ROWS):
        x0, y0, x1, y1 = cell(i)
        d.rounded_rectangle([x0, y0, x1, y1], radius=10,
                            fill=rgb(PANEL), outline=rgb(LINE), width=2)
        # faint code lines
        for j in range(2 + (i % 3)):
            ly = y0 + 16 + j * 14
            col = rgb(BLUE) if (i + j) % 4 == 1 else (rgb(VIOLET) if (i + j) % 4 == 2 else (43, 51, 64))
            d.rounded_rectangle([x0 + 12, ly, x0 + 12 + cw * (0.45 + ((i*13+j*7) % 40)/100), ly + 6],
                                radius=3, fill=col)

    # call-graph hints (sparse)
    for i in range(COLS * ROWS):
        c, r = i % COLS, i // COLS
        if c < COLS - 1 and (i * 7 + r) % 3 == 0:
            d.line([center(i), center(i + 1)], fill=(57, 66, 79, 180), width=2)
        if i + COLS < COLS * ROWS and (i * 5 + 1) % 3 == 0:
            d.line([center(i), center(i + COLS)], fill=(57, 66, 79, 180), width=2)

    # the bug path threading through components
    path = [2, 18, bug]
    pts = [center(i) for i in path]
    d.line(pts, fill=rgb(RED), width=5, joint="curve")
    for p in pts[:-1]:
        d.ellipse([p[0]-7, p[1]-7, p[0]+7, p[1]+7], fill=rgb(RED))
    bx0, by0, bx1, by1 = cell(bug)
    d.rounded_rectangle([bx0, by0, bx1, by1], radius=10, outline=rgb(RED), width=4)
    d.ellipse([bx1-30, by0+10, bx1-10, by0+30], fill=rgb(RED))

    # central play button
    cx, cy, rr = W // 2, H // 2 + 20, 92
    d.ellipse([cx-rr, cy-rr, cx+rr, cy+rr], fill=(5, 7, 11, 200), outline=(255, 255, 255, 230), width=4)
    tw = 46
    d.polygon([(cx-tw//2+6, cy-tw), (cx-tw//2+6, cy+tw), (cx+tw, cy)], fill=(255, 255, 255, 235))

    # brand
    try:
        f1 = ImageFont.truetype("DejaVuSans-Bold.ttf", 34)
        f2 = ImageFont.truetype("DejaVuSans.ttf", 24)
    except Exception:
        f1 = f2 = ImageFont.load_default()
    d.ellipse([60, 56, 84, 80], fill=rgb(VIOLET))
    d.text((96, 52), "bmc-agent", font=f1, fill=rgb(TEXT))
    d.text((W-470, 58), "agentic model checking", font=f2, fill=rgb(MUTED))
    img.save(POSTER)

make_poster()

# ===================================================================
# 2) build the slide
# ===================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = C(BG)

def no_shadow(shp):
    shp.shadow.inherit = False

def box(x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(shape)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = C(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = C(line); shape.line.width = Pt(line_w)
    return shape

def text(x, y, w, h, runs, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font="Calibri", spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, color, bold, size)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    for seg in runs:
        if seg == "\n":
            p = tf.add_paragraph(); p.alignment = align; p.line_spacing = spacing
            continue
        t, c, b, s = (seg + (None,)*4)[:4]
        r = p.add_run(); r.text = t
        r.font.size = Pt(s or size); r.font.bold = bool(b)
        r.font.color.rgb = C(c or color); r.font.name = font
    return tb

# ---- header / brand ----
dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(0.42), Inches(0.16), Inches(0.16))
no_shadow(dot); dot.fill.solid(); dot.fill.fore_color.rgb = C(VIOLET); dot.line.fill.background()
text(0.78, 0.32, 5, 0.4, [("agentic-prover", TEXT, True, 16),
                          ("   ·  agentic model checking", MUTED, False, 13)],
     font="Consolas")

# ---- title ----
text(0.55, 0.92, 8.6, 1.5,
     [("Who verifies the code ", TEXT, True, 38), ("\n", ), ("the AI just wrote?", BLUE, True, 38)],
     spacing=0.98)
text(0.57, 2.18, 8.6, 0.6,
     [("LLM-synthesized specs  ", MUTED, False, 14), ("×", VIOLET, True, 14),
      ("  bounded model checking  ", MUTED, False, 14), ("×", VIOLET, True, 14),
      ("  CEGAR refinement", MUTED, False, 14)])

# ---- pipeline (3 moves) ----
py, pw, ph = 2.78, 2.42, 1.18
nodes = [
    ("01 · synthesize", "LLM spec synthesis", "pre / post-conditions", VIOLET, "✦"),
    ("02 · check", "Bounded model check", "CBMC · Kani", BLUE, "⚙"),
    ("03 · classify", "Classify counterexample", "real bug vs false alarm", GREEN, "⚖"),
]
xs = [0.55, 0.55 + pw + 0.33, 0.55 + 2*(pw + 0.33)]
for (x, (step, title, sub, col, ico)) in zip(xs, nodes):
    box(x, py, pw, ph, fill=PANEL, line=col, line_w=1.5)
    text(x+0.18, py+0.12, pw-0.3, 0.3, [(ico+"  ", col, True, 16), (step, MUTED, False, 10.5)], font="Consolas")
    text(x+0.18, py+0.46, pw-0.3, 0.4, [(title, TEXT, True, 14.5)])
    text(x+0.18, py+0.82, pw-0.3, 0.3, [(sub, MUTED, False, 11)])
for ax in (xs[0]+pw+0.02, xs[1]+pw+0.02):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(py+ph/2-0.11), Inches(0.29), Inches(0.22))
    no_shadow(a); a.fill.solid(); a.fill.fore_color.rgb = C(BLUE); a.line.fill.background()
text(0.55, py+ph+0.06, 3*pw+0.66, 0.3,
     [("↻  CEGAR — refine the spec & re-check until the verdict is solid", AMBER, False, 12)],
     align=PP_ALIGN.CENTER)

# ---- composition point ----
cy0 = 4.62
text(0.55, cy0, 8.3, 0.4,
     [("You can’t check it part-by-part — ", TEXT, True, 14),
      ("the bug threads through the composition:", MUTED, False, 14)])
chips = [("parse_chunk", PANEL, LINE, TEXT), ("read_tag", PANEL, LINE, TEXT), ("buf[idx]", PANEL, RED, RED)]
cx = 0.55
for i, (label, fill, ln, tc) in enumerate(chips):
    w = 0.30 + 0.105 * len(label)
    box(cx, cy0+0.42, w, 0.42, fill=fill, line=ln, line_w=1.5)
    text(cx, cy0+0.42, w, 0.42, [(label, tc, True, 12.5)], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    cx += w
    if i < len(chips)-1:
        text(cx, cy0+0.42, 0.34, 0.42, [("→", MUTED, True, 15)], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE); cx += 0.34
text(0.55, cy0+0.98, 8.3, 0.4,
     [("spec ", MUTED, False, 11.5), ("idx < len", TEXT, True, 11.5), (" violated  ·  ", MUTED, False, 11.5),
      ("len=8, idx=12", TEXT, True, 11.5), ("  ·  replayed under ASan → ", MUTED, False, 11.5),
      ("SEGV", RED, True, 11.5)], font="Consolas")

# ---- confidence tiers (bottom strip) ----
ty = 6.42
tiers = [("confirmed_dynamic", GREEN), ("confirmed_system_entry", BLUE), ("confirmed_bmc", GREY)]
tx = 0.55
for label, col in tiers:
    w = 0.55 + 0.092 * len(label)
    box(tx, ty, w, 0.42, fill=PANEL, line=LINE, line_w=1.0)
    od = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(tx+0.14), Inches(ty+0.15), Inches(0.12), Inches(0.12))
    no_shadow(od); od.fill.solid(); od.fill.fore_color.rgb = C(col); od.line.fill.background()
    text(tx+0.34, ty, w-0.4, 0.42, [(label, TEXT, False, 11)], anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    tx += w + 0.18

# ===================================================================
# 3) demo panel (right) — embed movie if present, else poster + link
# ===================================================================
DX, DY, DW = 9.05, 2.62, 3.73
DH = DW * 9/16 + 0.34          # 16:9 + titlebar
# outer frame
box(DX-0.06, DY-0.06, DW+0.12, DH+0.12, fill=BG2, line=LINE, line_w=1.25)
# titlebar
box(DX, DY, DW, 0.30, fill=PANEL2, line=None)
for k, dc in enumerate(["ff5f56", "ffbd2e", "27c93f"]):
    od = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(DX+0.14+k*0.18), Inches(DY+0.10), Inches(0.10), Inches(0.10))
    no_shadow(od); od.fill.solid(); od.fill.fore_color.rgb = C(dc); od.line.fill.background()
text(DX+0.74, DY, DW-1.4, 0.30, [("bmc-agent demo", MUTED, False, 10)], anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
text(DX+DW-1.1, DY, 1.0, 0.30, [("● looping", GREEN, False, 9)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

vx, vy, vw, vh = DX, DY+0.30, DW, DW*9/16
if os.path.exists(MP4):
    mov = slide.shapes.add_movie(MP4, Inches(vx), Inches(vy), Inches(vw), Inches(vh),
                                 poster_frame_image=POSTER, mime_type="video/mp4")
    cap = "▶ Embedded video — plays in slideshow"
else:
    pic = slide.shapes.add_picture(POSTER, Inches(vx), Inches(vy), Inches(vw), Inches(vh))
    pic.click_action.hyperlink.address = "bmc-agent.html"   # open live demo on click
    cap = "▶ Click to play the live demo (bmc-agent.html)"
text(DX, vy+vh+0.06, DW, 0.3, [(cap, "5a6573", False, 10.5)], align=PP_ALIGN.CENTER)

# tagline bottom-right
text(DX-0.06, ty, DW+0.12, 0.42, [("bugs confirmed, not guessed.", TEXT, True, 13)],
     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("wrote", OUT, "(movie embedded)" if os.path.exists(MP4) else "(poster + hyperlink to HTML demo)")
